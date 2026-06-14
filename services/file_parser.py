"""
file_parser.py
Handles text extraction from PDF, PPTX, DOCX, and TXT files.
Cleans extracted text before returning.
Used by ai_service.py in Phase 2.
"""

import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import os
import sys


def clean_text(text: str) -> str:
    """Normalize raw extracted text by stripping noise and duplicate lines."""
    lines = text.splitlines()
    cleaned_lines = []
    seen_lower = set()

    for line in lines:
        stripped = line.strip()
        if stripped.isdigit():
            continue
        if len(stripped) < 3:
            continue
        lower = stripped.lower()
        if lower in seen_lower:
            continue
        seen_lower.add(lower)
        cleaned_lines.append(stripped)

    return "\n".join(cleaned_lines)


def extract_from_pdf(file_path: str) -> str:
    """Extract and clean text from all pages of a PDF file."""
    doc = fitz.open(file_path)
    parts = []
    for i in range(doc.page_count):
        page = doc.load_page(i)
        parts.append(page.get_text())
    doc.close()
    return clean_text("".join(parts))


def extract_from_pptx(file_path: str) -> str:
    """Extract and clean text from all slides in a PowerPoint file."""
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                parts.append(paragraph.text)
    return clean_text("\n".join(parts))


def extract_from_docx(file_path: str) -> str:
    """Extract and clean text from all paragraphs in a Word document."""
    doc = Document(file_path)
    parts = [paragraph.text for paragraph in doc.paragraphs]
    return clean_text("\n".join(parts))


def extract_from_txt(file_path: str) -> str:
    """Read and clean plain text from a UTF-8 text file."""
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        content = f.read()
    return clean_text(content)


def extract_text_from_file(file_path: str) -> str:
    """Route a file to the correct extractor based on its extension."""
    try:
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        if ext == ".pdf":
            return extract_from_pdf(file_path)
        if ext == ".pptx":
            return extract_from_pptx(file_path)
        if ext == ".docx":
            return extract_from_docx(file_path)
        if ext == ".txt":
            return extract_from_txt(file_path)

        print(f"[Parser Warning] Unsupported file type: {file_path}")
        return ""
    except Exception as e:
        print(f"[Parser Error] {file_path}: {e}")
        return ""


def extract_all(file_paths: list, pyq_flags: dict) -> dict:
    """Extract text from multiple files and split into study vs PYQ content."""
    study_parts = []
    pyq_parts = []

    for file_path in file_paths:
        text = extract_text_from_file(file_path)
        if not text:
            continue
        if pyq_flags.get(file_path, False):
            pyq_parts.append(text)
        else:
            study_parts.append(text)

    return {
        "study_content": "\n\n---FILE BREAK---\n\n".join(study_parts),
        "pyq_content": "\n\n---FILE BREAK---\n\n".join(pyq_parts),
    }


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python file_parser.py <file_path>")
        sys.exit(1)

    path = sys.argv[1]
    result = extract_text_from_file(path)

    if result:
        print("=== EXTRACTION SUCCESSFUL ===")
        print(f"Total characters: {len(result)}")
        print("\n--- First 500 characters ---")
        print(result[:500])
    else:
        print("=== EXTRACTION FAILED OR EMPTY ===")
